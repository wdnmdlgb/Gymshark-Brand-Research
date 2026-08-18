import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.util import Inches
import os

# 设置中文字体
plt.rcParams['font.sans-serif'] = ['SimHei', 'Microsoft YaHei', 'Arial Unicode MS']
plt.rcParams['axes.unicode_minus'] = False

base_dir = r"C:\Users\87090\Doubao\chats\2026-08-18\new-chat-1\项目1-Gymshark品牌研究"
img_dir = os.path.join(base_dir, "materials")
os.makedirs(img_dir, exist_ok=True)

# ========== 图表1：品牌能力雷达图 ==========
categories = ['品牌影响力', '产品性价比', '社区文化', 'KOL营销', '欧洲本地化', 'APP生态', '线下渠道']
N = len(categories)
angles = [n / float(N) * 2 * np.pi for n in range(N)]
angles += angles[:1]

gymshark = [8, 9, 10, 10, 7, 8, 6]
lululemon = [10, 4, 7, 6, 9, 7, 10]
alphalete = [5, 7, 6, 8, 3, 3, 4]
gymshark += gymshark[:1]
lululemon += lululemon[:1]
alphalete += alphalete[:1]

fig, ax = plt.subplots(figsize=(8, 8), subplot_kw=dict(polar=True))
ax.plot(angles, gymshark, 'o-', linewidth=2, label='Gymshark', color='#2F5496')
ax.fill(angles, gymshark, alpha=0.15, color='#2F5496')
ax.plot(angles, lululemon, 'o-', linewidth=2, label='Lululemon', color='#E8A87C')
ax.fill(angles, lululemon, alpha=0.12, color='#E8A87C')
ax.plot(angles, alphalete, 'o-', linewidth=2, label='Alphalete', color='#85CDCA')
ax.fill(angles, alphalete, alpha=0.12, color='#85CDCA')
ax.set_xticks(angles[:-1])
ax.set_xticklabels(categories, fontsize=12)
ax.set_ylim(0, 10)
ax.set_title('健身服饰品牌多维度能力对比', fontsize=15, fontweight='bold', pad=20)
ax.legend(loc='upper right', bbox_to_anchor=(1.3, 1.1), fontsize=11)
plt.tight_layout()
radar_path = os.path.join(img_dir, 'chart_radar.png')
plt.savefig(radar_path, dpi=150, bbox_inches='tight', facecolor='white')
plt.close()
print(f"雷达图已生成: {radar_path}")

# ========== 图表2：价格带柱状图 ==========
categories2 = ['Leggings', '训练短裤', '运动背心', '连帽衫', '运动内衣']
gymshark_prices = [65, 55, 35, 100, 42]
lululemon_prices = [120, 78, 68, 140, 60]
alphalete_prices = [70, 60, 45, 110, 50]

x = np.arange(len(categories2))
width = 0.25

fig, ax = plt.subplots(figsize=(10, 6))
bars1 = ax.bar(x - width, gymshark_prices, width, label='Gymshark', color='#2F5496')
bars2 = ax.bar(x, lululemon_prices, width, label='Lululemon', color='#E8A87C')
bars3 = ax.bar(x + width, alphalete_prices, width, label='Alphalete', color='#85CDCA')

ax.set_ylabel('价格 (USD)', fontsize=12)
ax.set_title('核心品类价格带对比（USD）', fontsize=15, fontweight='bold', pad=15)
ax.set_xticks(x)
ax.set_xticklabels(categories2, fontsize=11)
ax.legend(fontsize=11)
ax.grid(axis='y', alpha=0.3)

for bars in [bars1, bars2, bars3]:
    for bar in bars:
        height = bar.get_height()
        ax.annotate(f'${height}',
                    xy=(bar.get_x() + bar.get_width() / 2, height),
                    xytext=(0, 3), textcoords="offset points",
                    ha='center', va='bottom', fontsize=9)

plt.tight_layout()
bar_path = os.path.join(img_dir, 'chart_bar.png')
plt.savefig(bar_path, dpi=150, bbox_inches='tight', facecolor='white')
plt.close()
print(f"柱状图已生成: {bar_path}")

# ========== 图表3：用户痛点饼图 ==========
labels = ['尺码不一致', '面料耐用性', '客服与售后', '限量发售/APP体验', '运动内衣支撑不足']
sizes = [30, 25, 20, 15, 10]
colors = ['#2F5496', '#E8A87C', '#85CDCA', '#C38D9E', '#E27D60']
explode = (0.05, 0, 0, 0, 0)

fig, ax = plt.subplots(figsize=(9, 7))
wedges, texts, autotexts = ax.pie(sizes, explode=explode, labels=labels, colors=colors,
                                   autopct='%1.0f%%', shadow=False, startangle=90,
                                   textprops={'fontsize': 11})
for autotext in autotexts:
    autotext.set_color('white')
    autotext.set_fontweight('bold')
    autotext.set_fontsize(12)
ax.set_title('Gymshark用户负面反馈痛点分布', fontsize=15, fontweight='bold', pad=20)
plt.tight_layout()
pie_path = os.path.join(img_dir, 'chart_pie.png')
plt.savefig(pie_path, dpi=150, bbox_inches='tight', facecolor='white')
plt.close()
print(f"饼图已生成: {pie_path}")

# ========== 把图片插入PPT ==========
ppt_path = os.path.join(base_dir, "report", "Gymshark品牌深度研究报告.pptx")
prs = Presentation(ppt_path)

# 第4页（index=3）插入雷达图 - 替换右侧占位框
slide4 = prs.slides[3]
# 找到右侧的圆角矩形占位框并删除，然后插入图片
for shape in slide4.shapes:
    if shape.shape_type == 5 and shape.left > Inches(7):  # 圆角矩形在右侧
        sp = shape._element
        sp.getparent().remove(sp)
slide4.shapes.add_picture(radar_path, Inches(7.8), Inches(1.3), width=Inches(5.0))

# 第5页（index=4）插入柱状图
slide5 = prs.slides[4]
for shape in slide5.shapes:
    if shape.shape_type == 5 and shape.left > Inches(6.5):
        sp = shape._element
        sp.getparent().remove(sp)
slide5.shapes.add_picture(bar_path, Inches(6.8), Inches(1.3), width=Inches(6.0))

# 第11页（index=10）插入饼图
slide11 = prs.slides[10]
for shape in slide11.shapes:
    if shape.shape_type == 5 and shape.left < Inches(6):
        sp = shape._element
        sp.getparent().remove(sp)
slide11.shapes.add_picture(pie_path, Inches(0.4), Inches(1.8), width=Inches(5.8))

prs.save(ppt_path)
print(f"PPT已更新，图表已插入: {ppt_path}")

print("\n全部完成！")
