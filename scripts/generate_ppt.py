from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 配色方案
DARK_BLUE = RGBColor(0x1F, 0x3A, 0x5F)
MEDIUM_BLUE = RGBColor(0x2F, 0x54, 0x96)
LIGHT_BLUE = RGBColor(0xD6, 0xE4, 0xF0)
ACCENT_ORANGE = RGBColor(0xE8, 0xA8, 0x7C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)

def add_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=14, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT, font_name='微软雅黑'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_textbox(slide, left, top, width, height, items, font_size=12, color=DARK_GRAY, line_spacing=1.3):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = '微软雅黑'
        p.space_after = Pt(6)
        p.line_spacing = line_spacing
    return txBox

def add_title_bar(slide, title, subtitle=None):
    # 顶部色条
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    # 标题
    add_textbox(slide, 0.6, 0.2, 12, 0.7, title, font_size=24, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, 0.6, 0.65, 12, 0.4, subtitle, font_size=12, color=LIGHT_BLUE)

def add_table(slide, left, top, width, height, data, header_color=MEDIUM_BLUE):
    rows = len(data)
    cols = len(data[0])
    table = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)).table
    for i in range(rows):
        for j in range(cols):
            cell = table.cell(i, j)
            cell.text = str(data[i][j])
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                paragraph.font.name = '微软雅黑'
                if i == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                    paragraph.alignment = PP_ALIGN.CENTER
                else:
                    paragraph.font.color.rgb = DARK_GRAY
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
    return table

# ========== 第1页：封面 ==========
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide1, DARK_BLUE)
# 装饰线
shape = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2.8), Inches(3), Inches(0.08))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_ORANGE
shape.line.fill.background()

add_textbox(slide1, 1, 1.5, 11, 1.2, 'Gymshark 品牌深度研究', font_size=44, bold=True, color=WHITE)
add_textbox(slide1, 1, 3.1, 11, 0.8, '从车库创业到全球健身DTC标杆', font_size=24, color=LIGHT_BLUE)
add_textbox(slide1, 1, 4.0, 11, 0.6, '欧洲健身服饰市场竞品分析与国产健身品牌出海策略参考', font_size=16, color=LIGHT_BLUE)
add_textbox(slide1, 1, 5.8, 11, 0.4, '个人桌面研究项目  |  2026年8月', font_size=14, color=LIGHT_BLUE)
add_textbox(slide1, 1, 6.3, 11, 0.4, '研究对象：Gymshark（核心）| Lululemon | Alphalete（竞品）', font_size=12, color=LIGHT_BLUE)

# ========== 第2页：项目背景与研究框架 ==========
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide2, WHITE)
add_title_bar(slide2, '项目背景与研究框架', 'Project Background & Research Framework')

add_textbox(slide2, 0.6, 1.4, 6, 0.4, '研究背景', font_size=16, bold=True, color=MEDIUM_BLUE)
add_bullet_textbox(slide2, 0.6, 1.9, 6, 2.2, [
    '• 全球健身服饰市场持续增长，DTC品牌正在重塑传统运动服饰行业格局',
    '• Gymshark从2012年伯明翰车库起步，12年成长为年营收超6亿英镑的全球品牌',
    '• 欧洲作为全球第二大健身服饰市场，国产品牌出海面临本地化多重挑战',
    '• 拆解Gymshark成功逻辑，对中国健身服饰品牌出海具有重要参考价值'
], font_size=12)

add_textbox(slide2, 7, 1.4, 5.5, 0.4, '研究对象', font_size=16, bold=True, color=MEDIUM_BLUE)
data2 = [
    ['品牌', '国家', '定位'],
    ['Gymshark', '英国', '中端健身DTC+社区'],
    ['Lululemon', '加拿大', '高端瑜伽生活方式'],
    ['Alphalete', '美国', '中高端健身DTC+美学']
]
add_table(slide2, 7, 1.9, 5.5, 1.8, data2)

add_textbox(slide2, 0.6, 4.3, 12, 0.4, '8大研究维度', font_size=16, bold=True, color=MEDIUM_BLUE)
add_bullet_textbox(slide2, 0.6, 4.8, 12, 2, [
    '① 品牌基础与发展历程   ② 产品线与价格策略   ③ 爆款单品逻辑拆解   ④ 营销打法与KOL策略',
    '⑤ 标志性社区活动（Gymshark66）   ⑥ 欧洲本地化布局   ⑦ APP产品与品牌协同   ⑧ 用户反馈与口碑分析'
], font_size=13)

add_textbox(slide2, 0.6, 6.6, 12, 0.4, '数据来源：品牌官网、Companies House财报、App Store评价、Trustpilot、Morning Consult、Modash.io、行业媒体', font_size=10, color=RGBColor(0x88,0x88,0x88))

# ========== 第3页：品牌发展历程 ==========
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide3, WHITE)
add_title_bar(slide3, 'Gymshark品牌发展历程', 'Brand Development Milestones')

data3 = [
    ['年份', '关键事件', '意义'],
    ['2012', 'Ben Francis（19岁）伯明翰车库创立，缝纫机+丝网印刷', '品牌诞生，创始人IP起点'],
    ['2013', 'BodyPower健身展半小时销售额£30,000', '首次验证产品市场需求'],
    ['2015', 'Steve Hewitt出任CEO，专业化运营转型', '从个人作坊转向企业化'],
    ['2018', '首次外部融资，估值超£1B', '成为独角兽企业'],
    ['2020', '疫情居家健身红利，营收爆发增长', '加速全球化扩张'],
    ['2024', '财年营收£607.3m，连续12年增长', '规模突破6亿英镑'],
    ['2025.4', '阿姆斯特丹开设欧盟首家旗舰店', '欧洲线下渗透启动'],
    ['2026.2', '进入德国，Engelhorn+Breuninger店中店', '中欧市场拓展']
]
add_table(slide3, 0.6, 1.5, 12, 4.5, data3)

add_textbox(slide3, 0.6, 6.2, 12, 0.4, '创始人IP效应', font_size=16, bold=True, color=MEDIUM_BLUE)
add_bullet_textbox(slide3, 0.6, 6.6, 12, 0.8, [
    '• 19岁大学生车库创业的"草根逆袭"叙事，契合Z世代情感共鸣  • 个人社交媒体持续分享品牌幕后，建立真实感  • 创始人本身是健身爱好者，产品从自身需求出发'
], font_size=11)

# ========== 第4页：品牌基础横向对比 ==========
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide4, WHITE)
add_title_bar(slide4, '品牌基础横向对比', 'Competitive Benchmarking')

data4 = [
    ['对比维度', 'Gymshark', 'Lululemon', 'Alphalete'],
    ['创立时间', '2012年', '1998年', '2015年'],
    ['总部', '英国Solihull', '加拿大温哥华', '美国德州'],
    ['2024营收', '£607.3m($750m)', '约$10B', '估$200-300m'],
    ['估值', '约$1.45B', '约$50B+', '未公开'],
    ['核心定位', '健身DTC+社区文化', '高端瑜伽生活方式', '健身DTC+美学导向'],
    ['目标人群', '16-30岁，男女均衡', '25-45岁高收入女性', '18-30岁健身男性'],
    ['社媒粉丝', '1800万+', '约2000万+', '约500万+']
]
add_table(slide4, 0.6, 1.4, 7.5, 4.2, data4)

# 右侧雷达图占位
shape = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.4), Inches(1.4), Inches(4.3), Inches(4.2))
shape.fill.solid()
shape.fill.fore_color.rgb = LIGHT_GRAY
shape.line.color.rgb = MEDIUM_BLUE
add_textbox(slide4, 8.6, 1.6, 4, 0.4, '品牌多维度能力雷达图', font_size=13, bold=True, color=MEDIUM_BLUE, alignment=PP_ALIGN.CENTER)
add_textbox(slide4, 8.6, 3.0, 4, 1.5, '【图表占位】\n\nGymshark：社区文化10 / KOL营销10 / 性价比9\nLululemon：品牌影响力10 / 线下渠道10\nAlphalete：KOL营销8，其余维度规模较小', font_size=10, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide4, 0.6, 5.9, 12, 0.4, '关键发现', font_size=16, bold=True, color=MEDIUM_BLUE)
add_bullet_textbox(slide4, 0.6, 6.3, 12, 1, [
    '• Gymshark健身训练品类心智份额36%，超过Under Armour（35%），品牌认知度仅13%但垂直人群渗透率极高',
    '• 三家品牌形成"高端—中端—小众"三层定位，Gymshark在中端价格带建立差异化竞争优势'
], font_size=11)

# ========== 第5页：产品线与价格策略 ==========
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide5, WHITE)
add_title_bar(slide5, '产品线与价格策略', 'Product Line & Pricing Strategy')

data5 = [
    ['品类', 'Gymshark价格(USD)', 'Lululemon价格(USD)', 'Alphalete价格(USD)'],
    ['女款Leggings', '$55-$75', '$98-$148', '$60-$80'],
    ['男款训练短裤', '$45-$65', '$68-$88', '$50-$70'],
    ['运动背心', '$30-$40', '$58-$78', '$38-$52'],
    ['连帽衫', '$89-$115', '$118-$168', '$95-$130'],
    ['运动内衣', '$35-$50', '$52-$68', '$42-$58'],
    ['配饰', '$15-$45', '$38-$78', '$20-$50']
]
add_table(slide5, 0.6, 1.4, 6.5, 3.8, data5)

# 右侧柱状图占位
shape = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.4), Inches(1.4), Inches(5.3), Inches(3.8))
shape.fill.solid()
shape.fill.fore_color.rgb = LIGHT_GRAY
shape.line.color.rgb = MEDIUM_BLUE
add_textbox(slide5, 7.6, 1.6, 5, 0.4, '核心品类价格带对比柱状图', font_size=13, bold=True, color=MEDIUM_BLUE, alignment=PP_ALIGN.CENTER)
add_textbox(slide5, 7.6, 3.0, 5, 1.5, '【图表占位】\n\nGymshark全品类价格三家最低\nLeggings比Lululemon低约46%\n精准卡位$40-80中端价格带', font_size=11, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide5, 0.6, 5.5, 12, 0.4, '产品策略特征', font_size=16, bold=True, color=MEDIUM_BLUE)
add_bullet_textbox(slide5, 0.6, 5.9, 12, 1.4, [
    '• 高频上新+限量发售：潮牌drop模式制造稀缺感，Back In Stock页面持续引流',
    '• KOL联名系列：Whitney Simmons、David Laid联名，粉丝经济+限定发售',
    '• 男女均衡发展：不同于Lululemon（女性为主）和Alphalete（男性为主），女性bestsellers达175+款',
    '• 场景延伸：从纯训练服饰延伸到日常休闲、通勤穿搭，扩大使用场景'
], font_size=11)

# ========== 第6页：爆款单品逻辑拆解 ==========
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide6, WHITE)
add_title_bar(slide6, 'Gymshark爆款单品逻辑拆解', 'Best-Seller Product Analysis')

data6 = [
    ['爆款单品', '品类', '价格(USD)', '爆款原因'],
    ['Vital Seamless Leggings', '女款紧身裤', '$55-$65', '无缝针织+提臀剪裁，squat-proof口碑，社媒高频露出'],
    ['Power 5" Shorts(蓝标)', '男款训练短裤', '$55-$62', '标志性蓝标，健身圈文化符号，得物炒至¥1289'],
    ['Everyday Racer Back Tank', '女款运动背心', '$38', '基础百搭+内置胸垫，日常训练两用，复购率高'],
    ['Flex Training Shorts', '男款训练短裤', '$55-$65', '轻量化面料，适合HIIT/瑜伽，多场景穿着'],
    ['Whitney Simmons联名', '女款联名系列', '$60-$80', '网红联名+粉丝经济，限定发售制造稀缺']
]
add_table(slide6, 0.6, 1.4, 12, 3.2, data6)

add_textbox(slide6, 0.6, 4.9, 4, 0.4, '产品端', font_size=14, bold=True, color=MEDIUM_BLUE)
add_bullet_textbox(slide6, 0.6, 5.3, 4, 1.8, [
    '• 功能性优先：squat-proof、速干、高弹',
    '• 美学设计：修身剪裁凸显肌肉线条',
    '• 多色系选择：每季新配色刺激复购'
], font_size=11)

add_textbox(slide6, 4.8, 4.9, 4, 0.4, '营销端', font_size=14, bold=True, color=MEDIUM_BLUE)
add_bullet_textbox(slide6, 4.8, 5.3, 4, 1.8, [
    '• KOL共创：运动员参与产品设计',
    '• 社区验证：Gymshark Family内测',
    '• UGC裂变：用户穿搭帖自然传播'
], font_size=11)

add_textbox(slide6, 9, 4.9, 3.8, 0.4, '渠道端', font_size=14, bold=True, color=MEDIUM_BLUE)
add_bullet_textbox(slide6, 9, 5.3, 3.8, 1.8, [
    '• DTC独立站为主（96%销售）',
    '• 限量drop模式引发抢购',
    '• Back In Stock补货二次引流'
], font_size=11)

# ========== 第7页：营销打法对比分析 ==========
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide7, WHITE)
add_title_bar(slide7, '营销打法对比分析', 'Marketing Strategy Comparison')

data7 = [
    ['营销维度', 'Gymshark', 'Lululemon', 'Alphalete'],
    ['核心渠道', 'TikTok 69.2% + IG 30.5%', 'Instagram + 线下门店', 'YouTube + Instagram'],
    ['KOL策略', '微影响者矩阵+头部持股', '高端瑜伽导师长期合作', '创始人IP+健身网红'],
    ['内容风格', '真实、励志、训练日常', '生活方式、高端、正念', '硬核健身、美学展示'],
    ['付费模式', 'TikTok Spark Ads+Reels', '品牌广告+门店体验', 'YouTube长视频'],
    ['社区运营', 'Gymshark Family+大使', '门店瑜伽课+会员', 'YouTube社群']
]
add_table(slide7, 0.6, 1.4, 12, 2.8, data7)

add_textbox(slide7, 0.6, 4.5, 12, 0.4, 'Gymshark KOL营销三大特征', font_size=16, bold=True, color=MEDIUM_BLUE)
add_bullet_textbox(slide7, 0.6, 5.0, 12, 2.2, [
    '① 分层KOL矩阵：头部运动员（Chris Bumstead持股）+ 腰部创作者联盟营销 + 微影响者UGC自发传播',
    '② 联盟营销（Affiliate）机制：按销售佣金获利而非固定发帖费，激励创作者持续产出内容，ROI可控',
    '③ 平台内容差异化：TikTok（快速训练技巧+挑战，16-25岁）/ Instagram（美学穿搭种草）/ YouTube（长教程+纪录片）',
    '• 数据：KOL合作内容TikTok占69.2%（27.4k帖子），以短视频为主的投放结构使获客成本远低于传统品牌'
], font_size=12)

# ========== 第8页：Gymshark66挑战案例 ==========
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide8, WHITE)
add_title_bar(slide8, 'Gymshark66挑战案例拆解', 'Case Study: Gymshark66 Challenge')

# 左侧活动机制
add_textbox(slide8, 0.6, 1.4, 5.5, 0.4, '活动机制', font_size=16, bold=True, color=MEDIUM_BLUE)
add_bullet_textbox(slide8, 0.6, 1.9, 5.5, 2.5, [
    '• 周期：66天（行为科学：形成长期习惯的科学周期）',
    '• 规则：选目标→上传初始照→坚持66天→更新对比照',
    '• 激励：完成者有机会赢得一年Gymshark商品',
    '• 低门槛：不需要购买产品即可参与，最大化参与人数'
], font_size=12)

# 右侧传播效果
add_textbox(slide8, 7, 1.4, 5.5, 0.4, '传播效果', font_size=16, bold=True, color=MEDIUM_BLUE)
data8 = [
    ['指标', '数据'],
    ['Instagram帖子数', '110万+'],
    ['TikTok互动量', '6570万+'],
    ['跨平台蔓延', 'FB/Twitter/IG全平台'],
    ['举办频率', '每年固定，形成品牌传统']
]
add_table(slide8, 7, 1.9, 5.5, 2.2, data8)

add_textbox(slide8, 0.6, 4.7, 12, 0.4, '成功四要素', font_size=16, bold=True, color=MEDIUM_BLUE)
add_bullet_textbox(slide8, 0.6, 5.2, 12, 2, [
    '① 时机精准：新年后"健身决心"高峰期+放弃率最高时段，社区监督机制精准切入痛点',
    '② 社区属性：将营销转化为"community accountability ritual"（互相监督仪式），情感连接远超促销',
    '③ APP协同闭环：Training App提供免费训练计划→活动参与→APP打卡→社区分享→服饰购买→UGC传播',
    '④ UGC内容资产：110万+帖子构成巨大免费内容资产，真实用户故事比品牌广告更有说服力'
], font_size=12)

# ========== 第9页：欧洲本地化布局 ==========
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide9, WHITE)
add_title_bar(slide9, '欧洲本地化布局分析', 'European Localization Strategy')

add_textbox(slide9, 0.6, 1.4, 12, 0.4, '欧洲市场进入路径：先线上验证，后线下渗透', font_size=16, bold=True, color=MEDIUM_BLUE)
data9 = [
    ['阶段', '时间', '动作', '意义'],
    ['阶段1：纯DTC', '2012-2024', '14个欧洲本地化线上商店', '验证各国需求，积累用户数据'],
    ['阶段2：欧盟首店', '2025.4', '阿姆斯特丹Kalverstraat旗舰店', '欧盟线下品牌体验起点'],
    ['阶段3：中欧渗透', '2026.2', '德国Engelhorn(曼海姆)+Breuninger(斯图加特)店中店', '进入欧洲最大经济体，轻资产模式']
]
add_table(slide9, 0.6, 1.9, 12, 1.8, data9)

add_textbox(slide9, 0.6, 4.0, 6, 0.4, '本地化运营细节', font_size=16, bold=True, color=MEDIUM_BLUE)
add_bullet_textbox(slide9, 0.6, 4.5, 6, 2.5, [
    '• 多语言站点：英/德/法/西/意/荷/爱等14个本地化站点',
    '• 本地支付：Klarna分期（德/法/荷/北欧）、iDeal（荷兰）',
    '• 学生折扣：UNiDAYS覆盖欧洲多国，精准触达16-25岁人群',
    '• 欧洲仓储：本地仓储提升配送速度，降低物流和关税成本'
], font_size=12)

add_textbox(slide9, 7, 4.0, 5.5, 0.4, '竞品欧洲布局对比', font_size=16, bold=True, color=MEDIUM_BLUE)
data9b = [
    ['维度', 'Gymshark', 'Lululemon', 'Alphalete'],
    ['欧洲线下', '旗舰店+店中店', '50+直营门店', '无'],
    ['本地支付', 'Klarna+iDeal', '多币种', '仅PayPal'],
    ['本地语言', '7+语言', '多语言', '仅英文']
]
add_table(slide9, 7, 4.5, 5.5, 1.8, data9b)

# ========== 第10页：APP产品与品牌协同 ==========
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide10, WHITE)
add_title_bar(slide10, 'APP产品与品牌协同', 'APP Ecosystem & Brand Synergy')

add_textbox(slide10, 0.6, 1.4, 6, 0.4, '双APP策略', font_size=16, bold=True, color=MEDIUM_BLUE)
data10 = [
    ['APP', '模式', '核心功能', '评分'],
    ['Training App', '100%免费无广告', '450+课程、自定义计划、训练记录、计时器', '正面为主'],
    ['Shopping App', '免费购物工具', '限量发售推送、订单追踪、早期访问', '4.7/5']
]
add_table(slide10, 0.6, 1.9, 6, 1.6, data10)

add_textbox(slide10, 0.6, 3.8, 6, 0.4, '用户闭环逻辑', font_size=16, bold=True, color=MEDIUM_BLUE)
add_bullet_textbox(slide10, 0.6, 4.3, 6, 2.8, [
    '免费Training APP引流 → 用户注册开始训练',
    '↓',
    'Gymshark66活动增加粘性 → 训练中产生服饰需求',
    '↓',
    '跳转购物APP完成购买 → 穿着训练发布UGC',
    '↓',
    'UGC吸引新用户下载Training APP（循环）'
], font_size=12)

add_textbox(slide10, 7, 1.4, 5.5, 0.4, '竞品APP策略对比', font_size=16, bold=True, color=MEDIUM_BLUE)
data10b = [
    ['维度', 'Gymshark', 'Lululemon', 'Alphalete'],
    ['训练APP', '100%免费', '$39/月订阅', '无'],
    ['课程量', '450+', '10000+', 'YouTube'],
    ['硬件协同', '无', 'Mirror(已收缩)', '无'],
    ['协同逻辑', '免费引流→服饰', '内容订阅→硬件', 'YouTube→官网']
]
add_table(slide10, 7, 1.9, 5.5, 2.2, data10b)

add_textbox(slide10, 7, 4.4, 5.5, 0.4, '关键差异', font_size=16, bold=True, color=MEDIUM_BLUE)
add_bullet_textbox(slide10, 7, 4.9, 5.5, 2, [
    '• Lululemon走"高端订阅+硬件"路线（类Peloton），但Mirror已收缩',
    '• Gymshark完全免费策略降低门槛，用规模驱动服饰销售',
    '• 事实证明免费策略更利于大规模获客，契合中端品牌定位'
], font_size=11)

# ========== 第11页：用户反馈与口碑分析 ==========
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide11, WHITE)
add_title_bar(slide11, '用户反馈与口碑分析', 'User Feedback & Reputation Analysis')

add_textbox(slide11, 0.6, 1.4, 6, 0.4, '负面反馈痛点分布', font_size=16, bold=True, color=MEDIUM_BLUE)

# 左侧饼图占位
shape = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.9), Inches(5.5), Inches(3.2))
shape.fill.solid()
shape.fill.fore_color.rgb = LIGHT_GRAY
shape.line.color.rgb = MEDIUM_BLUE
add_textbox(slide11, 0.8, 2.1, 5, 0.4, '用户痛点饼图', font_size=13, bold=True, color=MEDIUM_BLUE, alignment=PP_ALIGN.CENTER)
add_textbox(slide11, 0.8, 3.2, 5, 1.8, '【图表占位】\n\n尺码不一致 30% | 面料耐用性 25%\n客服售后 20% | 限量发售/APP 15%\n运动内衣支撑 10%', font_size=11, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide11, 6.5, 1.4, 6.2, 0.4, '痛点详情', font_size=16, bold=True, color=MEDIUM_BLUE)
data11 = [
    ['痛点', '占比', '核心问题'],
    ['尺码不一致', '30%', '各系列标准不统一，洗后缩水，尺码表缺臀围数据'],
    ['面料耐用性', '25%', '深蹲透视，缝线开裂，起球褪色，面料偏薄'],
    ['客服售后', '20%', '退款换货繁琐，回复慢，账户无故冻结'],
    ['限量发售/APP', '15%', 'drop秒空难抢，发售时APP卡顿'],
    ['运动内衣', '10%', '大胸支撑不足，肩带滑落，五金件易损']
]
add_table(slide11, 6.5, 1.9, 6.2, 3.2, data11)

add_textbox(slide11, 0.6, 5.4, 12, 0.4, '正面评价 & 整体口碑', font_size=16, bold=True, color=MEDIUM_BLUE)
add_bullet_textbox(slide11, 0.6, 5.8, 12, 1.5, [
    '• 正面关键词：面料质感厚实舒适(35%) / 塑形效果好(25%) / 性价比高(15%) / 社区文化(10%)',
    '• Trustpilot约4.0/5星，75%用户给5星；BBB有41+投诉部分未解决',
    '• 整体"爱憎分明"：忠实用户热爱社区和性价比，不满用户集中投诉品控和客服',
    '• 这是中端DTC品牌典型成长烦恼：规模扩张后供应链品控和客服体系跟不上增长速度'
], font_size=11)

# ========== 第12页：结论与出海建议 ==========
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide12, WHITE)
add_title_bar(slide12, '结论与国产健身品牌出海建议', 'Conclusion & Go-to-Europe Recommendations')

add_textbox(slide12, 0.6, 1.4, 6, 0.4, 'Gymshark成功五要素', font_size=16, bold=True, color=MEDIUM_BLUE)
add_bullet_textbox(slide12, 0.6, 1.9, 6, 2.5, [
    '① 社区文化先行：先建社区再卖产品，卖的是身份认同',
    '② 微KOL矩阵+联盟营销：按效果付费，ROI可控',
    '③ 中端价格带卡位：$40-80区间差异化竞争',
    '④ 免费内容生态闭环：免费APP获客→活动留存→服饰变现',
    '⑤ DTC+精选线下渐进：线上掌握数据，线下平衡体验与成本'
], font_size=12)

add_textbox(slide12, 7, 1.4, 5.7, 0.4, '国产品牌出海六大建议', font_size=16, bold=True, color=MEDIUM_BLUE)
add_bullet_textbox(slide12, 7, 1.9, 5.7, 2.5, [
    '• 市场：优先荷兰/德国，电商渗透高健身文化成熟',
    '• 产品：$50-75中端价格带，统一尺码标准',
    '• 营销：TikTok微KOL矩阵，设计社区挑战活动',
    '• 内容：免费训练内容/APP建立用户粘性',
    '• 渠道：先DTC验证，后核心城市店中店',
    '• 本地化：多语言+Klarna支付+UNiDAYS学生折扣'
], font_size=12)

# 风险提示
shape = slide12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(4.7), Inches(12.1), Inches(2.2))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xFD, 0xF0, 0xE5)
shape.line.color.rgb = ACCENT_ORANGE
add_textbox(slide12, 0.8, 4.9, 11.5, 0.4, '⚠️ 风险提示', font_size=14, bold=True, color=RGBColor(0xC0, 0x60, 0x20))
add_bullet_textbox(slide12, 0.8, 5.4, 11.5, 1.4, [
    '• 欧盟合规：CE认证、REACH化学品法规、GDPR数据保护，合规成本不可忽视',
    '• 品控一致性：Gymshark教训表明规模扩张后品控下滑严重损害口碑，供应链是生命线',
    '• 社区壁垒：欧洲本土健身社区成熟，新品牌建立信任需要时间；客服体系必须提前搭建多语言支持'
], font_size=11)

# 保存
output_path = r"C:\Users\87090\Doubao\chats\2026-08-18\new-chat-1\项目1-Gymshark品牌研究\report\Gymshark品牌深度研究报告.pptx"
prs.save(output_path)
print(f"PPT已生成: {output_path}")
print(f"共{len(prs.slides)}页幻灯片")
